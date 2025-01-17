import openai
import requests
from pptx import Presentation
from pptx.util import Inches

# Set up your OpenAI API key
openai.api_key = 'your api key'  # Replace 'your-api-key' with your actual OpenAI API key

# Define the content for each slide
slides = [
    {
        "title": "Introduction to AI",
        "content": "Artificial Intelligence (AI) is the simulation of human intelligence in machines...",
        "image_prompt": "A futuristic robot interacting with humans"
    },
    {
        "title": "Applications of AI",
        "content": "AI has numerous applications including healthcare, finance, transportation, and more...",
        "image_prompt": "AI in healthcare, finance, and transportation"
    },
    # Add more slides as needed
]

# Generate images with DALL·E 3
for slide in slides:
    response = openai.Image.create(
        prompt=slide["image_prompt"],
        n=1,
        size="1024x1024"
    )
    slide["image_url"] = response['data'][0]['url']

# Download images locally
for slide in slides:
    image_url = slide["image_url"]
    image_data = requests.get(image_url).content
    image_path = f'{slide["title"].replace(" ", "_")}.png'
    with open(image_path, 'wb') as handler:
        handler.write(image_data)
    slide["image_path"] = image_path

# Generate enhanced content with the Chat API
for slide in slides:
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Create a detailed explanation for the slide titled '{slide['title']}' with the content: {slide['content']}"}
        ]
    )
    slide["enhanced_content"] = response['choices'][0]['message']['content']

# Create a PowerPoint presentation
prs = Presentation()

for slide in slides:
    # Add a slide
    slide_layout = prs.slide_layouts[5]  # Choosing a blank layout
    slide_obj = prs.slides.add_slide(slide_layout)
    
    # Add title
    title = slide_obj.shapes.title
    title.text = slide["title"]
    
    # Add content
    textbox = slide_obj.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(4))
    text_frame = textbox.text_frame
    text_frame.text = slide["enhanced_content"]
    
    # Add image
    image_path = slide["image_path"]
    slide_obj.shapes.add_picture(image_path, Inches(5), Inches(1.5), width=Inches(4.5))

# Save the presentation
prs.save('AI_Presentation.pptx')

print("Presentation created successfully!")

